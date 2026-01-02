from __future__ import annotations


# =============================================================================
# Teaching notes: pptx_export.py (Editable PowerPoint export)
#
# This module exports the roadmap to a .pptx slide using python-pptx.
#
# Important difference from renderer.py:
#   - renderer.py produces pixels/vectors (PNG/PDF).
#   - pptx_export.py produces editable PowerPoint shapes.
#
# PowerPoint uses inches for positioning in python-pptx.
# So we convert "data coordinates" (dates/lanes) into slide inches.
# =============================================================================

"""Editable PPTX export.

This module renders the roadmap as editable PowerPoint shapes:
  - Workstreams as swimlanes
  - Tasks as rounded rectangles
  - Milestones as diamonds
  - Timeline header band (auto weeks/months/quarters/years)

The intent is an executive-ready slide that remains easy to tweak in
PowerPoint (move/resize blocks, edit text, etc.).
"""

# ---------------------------------------------------------------------------
# Imports (standard library, third-party packages, and local modules)
# ---------------------------------------------------------------------------
from dataclasses import dataclass
from datetime import date, datetime, timedelta
from io import BytesIO
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from zoneinfo import ZoneInfo

from date_utils import block_span_inclusive, date_to_x
from renderer import (
    DEFAULT_PALETTE,
    choose_timeline_mode,
    compute_bands_and_rows,
    resolve_font_family,
)
from roadmap_models import Settings, Task, Workstream
from scheduler import schedule_by_workstream


# ----------------------------
# Color helpers
# ----------------------------


def _hex_to_rgb(hex_color: str) -> RGBColor:
    s = (hex_color or "").strip()
    if not s:
        return RGBColor(0, 0, 0)
    if s.startswith("#"):
        s = s[1:]
    s = s.upper()
    r = int(s[0:2], 16)
    g = int(s[2:4], 16)
    b = int(s[4:6], 16)
    return RGBColor(r, g, b)


def _lighten_hex(hex_color: str, amount: float) -> str:
    """Lighten a #RRGGBB color by mixing with white.

    amount: 0 -> unchanged, 1 -> white
    """
    s = (hex_color or "").strip()
    if not s:
        return "#FFFFFF"
    if s.startswith("#"):
        s = s[1:]
    r = int(s[0:2], 16)
    g = int(s[2:4], 16)
    b = int(s[4:6], 16)
    r2 = int(round(r + (255 - r) * amount))
    g2 = int(round(g + (255 - g) * amount))
    b2 = int(round(b + (255 - b) * amount))
    return f"#{r2:02X}{g2:02X}{b2:02X}"


def _pick_workstream_colors(workstreams: List[Workstream]) -> Dict[str, str]:
    out: Dict[str, str] = {}
    i = 0
    for ws in workstreams:
        if ws.color:
            out[ws.workstream] = ws.color
        else:
            out[ws.workstream] = DEFAULT_PALETTE[i % len(DEFAULT_PALETTE)]
            i += 1
    return out


# ----------------------------
# Timeline helpers
# (importing renderer internals keeps PPTX output consistent with PDF/PNG)
# ----------------------------


from renderer import (  # noqa: E402  (intentional internal import)
    _build_month_segments,
    _build_quarter_segments,
    _build_week_segments,
    _build_year_segments,
    _iter_month_starts,
    _iter_quarter_starts,
    _iter_week_starts,
    _iter_year_starts,
)


# ----------------------------
# Text fitting (simple, deterministic)
# ----------------------------


def _wrap_and_truncate(text: str, width: int, max_lines: int) -> List[str]:
    import textwrap

    if max_lines <= 0:
        return []
    text = (text or "").strip()
    if not text:
        return [""]
    lines = textwrap.wrap(text, width=width) or [""]
    if len(lines) <= max_lines:
        return lines
    kept = lines[:max_lines]
    kept[-1] = _ellipsis(kept[-1], width)
    return kept


def _ellipsis(line: str, width: int) -> str:
    if width <= 1:
        return "…"
    line = (line or "").rstrip()
    if len(line) <= width:
        return line
    if width <= 2:
        return line[: width - 1] + "…"
    return line[: width - 1].rstrip() + "…"


def _is_truncated(lines: List[str], original: str) -> bool:
    if not lines:
        return False
    return lines[-1].endswith("…") and not (original or "").endswith("…")


@dataclass(frozen=True)
class FittedText:
    title_lines: List[str]
    desc_lines: List[str]
    font_size_pt: int


def fit_text_ppt(
    title: str,
    desc: Optional[str],
    *,
    width_in: float,
    height_in: float,
    preview: bool = False,
) -> FittedText:
    """Best-effort text fitting for PPT shapes.

    Uses the same heuristic as the matplotlib renderer (chars/line and max
    lines), expressed in points (72 pts per inch).
    """
    base = 12 if not preview else 11
    min_fs = 8 if not preview else 7

    width_pt = max(width_in * 72.0, 36.0)
    height_pt = max(height_in * 72.0, 18.0)

    title = (title or "").strip()
    desc = (desc or "").strip() if desc else ""

    best = FittedText([title], [], min_fs)

    for fs in range(base, min_fs - 1, -1):
        max_cpl = max(int(width_pt / (fs * 0.55)), 8)
        max_lines = max(int(height_pt / (fs * 1.25)), 1)

        title_max_lines = min(2, max_lines) if desc else max_lines
        title_lines = _wrap_and_truncate(title, max_cpl, title_max_lines)

        remaining = max_lines - len(title_lines)
        desc_lines: List[str] = []
        if desc and remaining > 0:
            desc_lines = _wrap_and_truncate(desc, max_cpl, remaining)

        fitted = FittedText(title_lines, desc_lines, fs)

        # Prefer first fontsize where title isn't truncated.
        if not _is_truncated(title_lines, title):
            return fitted
        best = fitted

    return best


# ----------------------------
# PPTX rendering
# ----------------------------


def export_pptx_bytes(
    settings: Settings,
    workstreams: List[Workstream],
    tasks: List[Task],
    *,
    include_out_of_range: bool = False,
) -> bytes:
    """Return an editable PPTX as bytes."""


    # Export process (high level):
    # 1) Make sure every task has a sublane (so blocks won't overlap).
    # 2) Normalize colors and clamp tasks to the visible date range.
    # 3) Create a blank slide sized like the PDF/PNG output (A3/A4 landscape).
    # 4) Compute lane geometry and convert dates → inches on the slide.
    # 5) Draw header labels, lane bands, blocks, milestones, and status styling.
    # 6) Save the PPTX into memory and return the raw bytes.
    # Ensure sublanes exist (robust for callers that forget to schedule).
    if any(t.sublane is None for t in tasks):
        scheduled = schedule_by_workstream(tasks, touching_counts_as_overlap=True)
        by_id = {t.id: t for ws_tasks in scheduled.values() for t in ws_tasks}
        tasks = [by_id[t.id] for t in tasks if t.id in by_id]

    # Resolve font
    font_name = resolve_font_family(settings.font_family)

    overall_start = settings.overall_start_date
    overall_end = settings.overall_end_date

    # Normalize workstream colors
    ws_color_map = _pick_workstream_colors(workstreams)
    workstreams_norm = [ws.model_copy(update={"color": ws_color_map[ws.workstream]}) for ws in workstreams]

    # Filter/clamp tasks like the matplotlib renderer
    visible_tasks: List[Task] = []
    for t in tasks:
        if t.end_date < overall_start or t.start_date > overall_end:
            if include_out_of_range:
                visible_tasks.append(t)
            continue

        clamped_start = max(t.start_date, overall_start)
        clamped_end = min(t.end_date, overall_end)
        if clamped_start != t.start_date or clamped_end != t.end_date:
            visible_tasks.append(t.model_copy(update={"start_date": clamped_start, "end_date": clamped_end}))
        else:
            visible_tasks.append(t)

    # Slide size: match PDF/PNG aspect (A3/A4 landscape) so exports align.
    if settings.page_size == "A4":
        slide_w_in, slide_h_in = 11.69, 8.27
    else:
        slide_w_in, slide_h_in = 16.54, 11.69

    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Layout regions
    margin_l, margin_r = 0.55, 0.45
    margin_t, margin_b = 0.45, 0.45

    usable_w = slide_w_in - margin_l - margin_r
    usable_h = slide_h_in - margin_t - margin_b

    header_h = usable_h * 0.13
    chart_h = usable_h - header_h

    label_w = usable_w * 0.23
    main_w = usable_w - label_w

    header_left = margin_l
    header_top = margin_t
    chart_left = margin_l
    chart_top = margin_t + header_h

    label_left = chart_left
    main_left = chart_left + label_w

    # Compute y layout in "units" and map to inches (same as matplotlib)
    timeline_mode = choose_timeline_mode(overall_start, overall_end)
    timeline_rows = 2 if timeline_mode in ("weeks", "quarters_years") else 1
    timeline_row_h_units = 0.65
    timeline_h_units = timeline_rows * timeline_row_h_units

    bands, row_map, total_height_units = compute_bands_and_rows(workstreams_norm, visible_tasks)

    total_units = total_height_units + timeline_h_units
    inch_per_unit = chart_h / max(total_units, 1e-6)

    timeline_row_h = timeline_row_h_units * inch_per_unit
    timeline_h = timeline_h_units * inch_per_unit

    lane_top = chart_top + timeline_h

    # X mapping helpers
    x0 = date_to_x(overall_start)
    x1 = date_to_x(overall_end + timedelta(days=1))
    total_x = max(x1 - x0, 1.0)

    def x_to_in(x: float) -> float:
        frac = (x - x0) / total_x
        return main_left + frac * main_w

    # Colors
    c_border = "#DADADA"
    c_grid = "#E6E6E6"
    c_major = "#D0D0D0"
    c_year = "#C8C8C8"
    c_label_bg = "#F6F8FB"
    c_band_alt = "#FAFAFA"
    c_text = "#222222"

    # -------------------------
    # Header
    # -------------------------
    title_box = slide.shapes.add_textbox(Inches(header_left), Inches(header_top), Inches(usable_w * 0.72), Inches(header_h * 0.55))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = settings.chart_title
    p.font.name = font_name
    p.font.bold = True
    p.font.size = Pt(30 if settings.page_size == "A3" else 24)
    p.font.color.rgb = _hex_to_rgb("#111827")

    if settings.chart_subtitle:
        sub_box = slide.shapes.add_textbox(Inches(header_left), Inches(header_top + header_h * 0.52), Inches(usable_w * 0.72), Inches(header_h * 0.35))
        tf2 = sub_box.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = settings.chart_subtitle
        p2.font.name = font_name
        p2.font.size = Pt(14 if settings.page_size == "A3" else 12)
        p2.font.color.rgb = _hex_to_rgb("#374151")

    if settings.confidentiality_label:
        conf_box = slide.shapes.add_textbox(Inches(margin_l + usable_w * 0.74), Inches(header_top), Inches(usable_w * 0.26), Inches(header_h * 0.35))
        tfc = conf_box.text_frame
        tfc.clear()
        pc = tfc.paragraphs[0]
        pc.text = settings.confidentiality_label
        pc.font.name = font_name
        pc.font.size = Pt(12)
        pc.font.color.rgb = _hex_to_rgb("#374151")
        pc.alignment = PP_ALIGN.RIGHT

    date_range_text = f"{overall_start.strftime('%d %b %Y')} – {overall_end.strftime('%d %b %Y')}"
    dr_box = slide.shapes.add_textbox(Inches(header_left), Inches(header_top + header_h * 0.82), Inches(usable_w * 0.6), Inches(header_h * 0.25))
    tfd = dr_box.text_frame
    tfd.clear()
    pd = tfd.paragraphs[0]
    pd.text = date_range_text
    pd.font.name = font_name
    pd.font.size = Pt(12)
    pd.font.color.rgb = _hex_to_rgb("#374151")

    # Divider line under header
    div = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(margin_l),
        Inches(chart_top),
        Inches(margin_l + usable_w),
        Inches(chart_top),
    )
    div.line.color.rgb = _hex_to_rgb("#E6E6E6")
    div.line.width = Pt(1.0)

    # -------------------------
    # Background panels
    # -------------------------
    # Label column background
    lab_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(label_left),
        Inches(chart_top),
        Inches(label_w),
        Inches(chart_h),
    )
    lab_bg.fill.solid()
    lab_bg.fill.fore_color.rgb = _hex_to_rgb(c_label_bg)
    lab_bg.line.fill.background()  # no outline

    # Main background (white)
    main_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(main_left),
        Inches(chart_top),
        Inches(main_w),
        Inches(chart_h),
    )
    main_bg.fill.solid()
    main_bg.fill.fore_color.rgb = _hex_to_rgb("#FFFFFF")
    main_bg.line.fill.background()

    # Vertical divider between labels and main
    vdiv = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(main_left),
        Inches(chart_top),
        Inches(main_left),
        Inches(chart_top + chart_h),
    )
    vdiv.line.color.rgb = _hex_to_rgb("#E0E0E0")
    vdiv.line.width = Pt(1.0)

    # Separator line between timeline band and swimlanes
    sep = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(main_left),
        Inches(lane_top),
        Inches(main_left + main_w),
        Inches(lane_top),
    )
    sep.line.color.rgb = _hex_to_rgb(c_major)
    sep.line.width = Pt(1.3)

    # Label side timeline separator as well
    sep2 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(label_left),
        Inches(lane_top),
        Inches(label_left + label_w),
        Inches(lane_top),
    )
    sep2.line.color.rgb = _hex_to_rgb(c_major)
    sep2.line.width = Pt(1.3)

    # -------------------------
    # Timeline header band
    # -------------------------

    def _add_timeline_row(
        *,
        row_index: int,
        kind: str,
        segs: List[Tuple[date, date, str]],
    ) -> None:
        top_in = chart_top + row_index * timeline_row_h
        for i, (ds, de, label) in enumerate(segs):
            xs = x_to_in(date_to_x(ds))
            xe = x_to_in(date_to_x(de))
            w_in = max(xe - xs, 0.02)
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(xs),
                Inches(top_in),
                Inches(w_in),
                Inches(timeline_row_h),
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = _hex_to_rgb("#FFFFFF" if (i % 2 == 0) else "#F7F7F7")
            rect.line.color.rgb = _hex_to_rgb(c_border)
            rect.line.width = Pt(0.8)

            # Centered label
            tb = slide.shapes.add_textbox(
                Inches(xs),
                Inches(top_in + timeline_row_h * 0.05),
                Inches(w_in),
                Inches(timeline_row_h * 0.9),
            )
            tf = tb.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.text = label
            p.font.name = font_name
            if kind == "weeks":
                p.font.size = Pt(10)
            elif kind == "months":
                p.font.size = Pt(11)
            elif kind == "quarters":
                p.font.size = Pt(12)
            else:
                p.font.size = Pt(12)
            p.font.color.rgb = _hex_to_rgb("#333333")

    if timeline_mode == "weeks":
        month_segs = _build_month_segments(overall_start, overall_end)
        week_segs = _build_week_segments(overall_start, overall_end, settings.week_start_day)
        _add_timeline_row(row_index=0, kind="months", segs=month_segs)
        _add_timeline_row(row_index=1, kind="weeks", segs=week_segs)
    elif timeline_mode == "months":
        month_segs = _build_month_segments(overall_start, overall_end)
        _add_timeline_row(row_index=0, kind="months", segs=month_segs)
    elif timeline_mode == "quarters":
        q_segs = _build_quarter_segments(overall_start, overall_end, include_year=True)
        _add_timeline_row(row_index=0, kind="quarters", segs=q_segs)
    else:
        y_segs = _build_year_segments(overall_start, overall_end)
        q_segs = _build_quarter_segments(overall_start, overall_end, include_year=False)
        _add_timeline_row(row_index=0, kind="years", segs=y_segs)
        _add_timeline_row(row_index=1, kind="quarters", segs=q_segs)

    # -------------------------
    # Vertical grid lines
    # -------------------------
    lane_bottom = chart_top + chart_h

    def _vline(d: date, *, color: str, width_pt: float) -> None:
        x = x_to_in(date_to_x(d))
        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x),
            Inches(lane_top),
            Inches(x),
            Inches(lane_bottom),
        )
        ln.line.color.rgb = _hex_to_rgb(color)
        ln.line.width = Pt(width_pt)

    if timeline_mode == "weeks":
        for d in _iter_week_starts(overall_start, overall_end, settings.week_start_day):
            _vline(d, color=c_grid, width_pt=0.75)
        for d in _iter_month_starts(overall_start, overall_end):
            _vline(d, color=c_major, width_pt=1.2)
    elif timeline_mode == "months":
        for d in _iter_month_starts(overall_start, overall_end):
            _vline(d, color=c_grid, width_pt=0.9)
        for d in _iter_year_starts(overall_start, overall_end):
            _vline(d, color=c_year, width_pt=1.4)
    else:
        for d in _iter_quarter_starts(overall_start, overall_end):
            _vline(d, color=c_grid, width_pt=1.0)
        for d in _iter_year_starts(overall_start, overall_end):
            _vline(d, color=c_year, width_pt=1.6)

    # -------------------------
    # Workstream bands + labels
    # -------------------------
    for i, band in enumerate(bands):
        y0_in = lane_top + band.y0 * inch_per_unit
        y1_in = lane_top + band.y1 * inch_per_unit
        h_in = max(y1_in - y0_in, 0.02)

        # alternating band fill
        if i % 2 == 0:
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(main_left),
                Inches(y0_in),
                Inches(main_w),
                Inches(h_in),
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = _hex_to_rgb(c_band_alt)
            rect.line.fill.background()

        # band separators
        for y in (y0_in, y1_in):
            ln = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(main_left),
                Inches(y),
                Inches(main_left + main_w),
                Inches(y),
            )
            ln.line.color.rgb = _hex_to_rgb(c_major)
            ln.line.width = Pt(1.0)

            ln2 = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(label_left),
                Inches(y),
                Inches(label_left + label_w),
                Inches(y),
            )
            ln2.line.color.rgb = _hex_to_rgb(c_major)
            ln2.line.width = Pt(1.0)

        # workstream accent bar + name
        ws_color = ws_color_map.get(band.workstream, "#1F77B4")
        accent_w = max(label_w * 0.06, 0.08)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(label_left + 0.12),
            Inches(y0_in + 0.06),
            Inches(accent_w),
            Inches(max(h_in - 0.12, 0.02)),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = _hex_to_rgb(ws_color)
        accent.line.fill.background()

        name_box = slide.shapes.add_textbox(
            Inches(label_left + 0.12 + accent_w + 0.12),
            Inches(y0_in),
            Inches(label_w - (0.12 + accent_w + 0.22)),
            Inches(h_in),
        )
        tf = name_box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = band.workstream
        p.font.name = font_name
        p.font.bold = True
        p.font.size = Pt(14 if settings.page_size == "A3" else 12)
        p.font.color.rgb = _hex_to_rgb(c_text)

    # sublane lines
    for (ws, lane), row in row_map.items():
        y = lane_top + row.y0 * inch_per_unit
        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(main_left),
            Inches(y),
            Inches(main_left + main_w),
            Inches(y),
        )
        ln.line.color.rgb = _hex_to_rgb("#EFEFEF")
        ln.line.width = Pt(0.6)

    # -------------------------
    # Today line
    # -------------------------
    if settings.show_today_line:
        tz = ZoneInfo(settings.timezone)
        today = settings.today_line_date
        if today is None:
            today = datetime.now(tz=tz).date()
        if overall_start <= today <= overall_end:
            xt = x_to_in(date_to_x(today))
            ln = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(xt),
                Inches(chart_top),
                Inches(xt),
                Inches(lane_bottom),
            )
            ln.line.color.rgb = _hex_to_rgb("#111111")
            ln.line.width = Pt(1.25)
            ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH

            # small label
            tb = slide.shapes.add_textbox(
                Inches(min(xt + 0.06, main_left + main_w - 0.6)),
                Inches(chart_top + 0.02),
                Inches(0.6),
                Inches(0.25),
            )
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "Today"
            p.font.name = font_name
            p.font.size = Pt(10)
            p.font.color.rgb = _hex_to_rgb("#111111")

    # -------------------------
    # Tasks
    # -------------------------
    # Status styling (clearer differentiation)
    # We match the PNG/PDF renderer: a left status stripe + stronger border styles.
    STATUS_STYLE = {
        "planned": {"stripe": "#6B7280", "edge": "#3A3A3A", "lw": 1.0, "dash": None, "lighten": 0.0},
        "in_progress": {"stripe": "#2563EB", "edge": "#2563EB", "lw": 2.0, "dash": MSO_LINE_DASH_STYLE.DASH, "lighten": 0.0},
        "done": {"stripe": "#16A34A", "edge": "#6B7280", "lw": 1.0, "dash": None, "lighten": 0.65},
        "risk": {"stripe": "#DC2626", "edge": "#DC2626", "lw": 2.25, "dash": None, "lighten": 0.0},
    }

    row_padding_units = 0.12
    for t in visible_tasks:
        lane = int(t.sublane or 0)
        row = row_map.get((t.workstream, lane))
        if row is None:
            continue

        # vertical placement
        y0_in = lane_top + (row.y0 + row_padding_units) * inch_per_unit
        lane_h_in = 1.0 * inch_per_unit
        pad_in = row_padding_units * inch_per_unit
        h_in = max(lane_h_in - 2 * pad_in, 0.05)

        ws_color = ws_color_map.get(t.workstream, "#1F77B4")
        face = t.color_override or ws_color

        status = (t.status or "planned").lower() if t.status else "planned"
        style = STATUS_STYLE.get(status, STATUS_STYLE["planned"])
        edge = style["edge"]
        lw = float(style["lw"])
        dash = style["dash"]
        stripe_color = style["stripe"]

        if float(style.get("lighten", 0.0)) > 0:
            face = _lighten_hex(face, float(style["lighten"]))

        if t.type == "milestone":
            cx = x_to_in(date_to_x(t.start_date))
            day_in = main_w / total_x
            diamond_w = max(day_in * 0.70, 0.12)
            diamond_h = min(h_in * 0.85, 0.32)
            left = cx - diamond_w / 2.0
            top = y0_in + (h_in - diamond_h) / 2.0
            shp = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                Inches(left),
                Inches(top),
                Inches(diamond_w),
                Inches(diamond_h),
            )
            shp.fill.solid()
            shp.fill.fore_color.rgb = _hex_to_rgb(face)
            shp.line.color.rgb = _hex_to_rgb(edge)
            shp.line.width = Pt(lw)
            if dash is not None:
                shp.line.dash_style = dash
            if t.hyperlink:
                shp.click_action.hyperlink.address = t.hyperlink

            # label to the right
            label_left = min(cx + diamond_w / 2.0 + 0.08, main_left + main_w - 0.2)
            label_w = max(main_left + main_w - label_left - 0.05, 0.4)
            tb = slide.shapes.add_textbox(
                Inches(label_left),
                Inches(y0_in),
                Inches(label_w),
                Inches(h_in),
            )
            ft = tb.text_frame
            ft.clear()
            ft.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            p = ft.paragraphs[0]
            p.text = t.title
            p.font.name = font_name
            p.font.size = Pt(12)
            p.font.color.rgb = _hex_to_rgb("#111111")
            continue

        # block
        bx0, bx1 = block_span_inclusive(t.start_date, t.end_date)
        if (bx1 - bx0) < 0.8:
            bx1 = bx0 + 0.8

        left = x_to_in(bx0)
        right = x_to_in(bx1)
        w_in = max(right - left, 0.08)

        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(y0_in),
            Inches(w_in),
            Inches(h_in),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = _hex_to_rgb(face)
        shp.line.color.rgb = _hex_to_rgb(edge)
        shp.line.width = Pt(lw)
        if dash is not None:
            shp.line.dash_style = dash

        if t.hyperlink:
            shp.click_action.hyperlink.address = t.hyperlink

        # Status stripe (kept thin and consistent)
        stripe_w = min(max(w_in * 0.12, 0.06), w_in * 0.35, 0.12)
        if stripe_w > 0:
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left),
                Inches(y0_in),
                Inches(stripe_w),
                Inches(h_in),
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = _hex_to_rgb(stripe_color)
            stripe.line.fill.background()
            if t.hyperlink:
                stripe.click_action.hyperlink.address = t.hyperlink

        # Text
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.06 + stripe_w)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)

        fitted = fit_text_ppt(
            t.title,
            t.description,
            width_in=max(w_in - (0.12 + stripe_w), 0.2),
            height_in=max(h_in - 0.06, 0.15),
        )

        # Title lines (bold)
        first = True
        for line in fitted.title_lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = line
            p.font.name = font_name
            p.font.bold = True
            p.font.size = Pt(fitted.font_size_pt)
            p.font.color.rgb = _hex_to_rgb("#1A1A1A" if status != "done" else "#4A4A4A")
            p.space_after = Pt(0)

        # Description lines
        for line in fitted.desc_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.name = font_name
            p.font.bold = False
            p.font.size = Pt(max(fitted.font_size_pt - 1, 7))
            p.font.color.rgb = _hex_to_rgb("#1A1A1A" if status != "done" else "#4A4A4A")
            p.space_after = Pt(0)

    # -------------------------
    # Minimal legend (only if it adds value)
    # -------------------------
    statuses_present = {((t.status or "planned").lower()) for t in visible_tasks}
    if len(statuses_present) > 1:
        order = {"planned": 0, "in_progress": 1, "risk": 2, "done": 3}
        statuses_sorted = sorted(statuses_present, key=lambda s: order.get(s, 99))
        legend_left = main_left
        legend_top = chart_top + chart_h + 0.05 - margin_b
        # If we don't have room below, tuck into header bottom-right.
        if legend_top + 0.35 > slide_h_in - 0.05:
            legend_top = header_top + header_h * 0.78
            legend_left = margin_l + usable_w * 0.62

        tb = slide.shapes.add_textbox(Inches(legend_left), Inches(legend_top), Inches(0.8), Inches(0.25))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Legend:"
        p.font.name = font_name
        p.font.size = Pt(11)
        p.font.color.rgb = _hex_to_rgb("#111111")

        label_map = {
            "planned": "Planned",
            "in_progress": "In progress",
            "done": "Done",
            "risk": "Risk",
        }

        x_cursor = legend_left + 0.85
        sample_w = 0.28
        sample_h = 0.18
        stripe_frac = 0.28

        for s in statuses_sorted:
            style = STATUS_STYLE.get(s, STATUS_STYLE["planned"])
            edge = style["edge"]
            lw = float(style["lw"])
            dash = style.get("dash")
            stripe_color = style["stripe"]

            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_cursor),
                Inches(legend_top + 0.02),
                Inches(sample_w),
                Inches(sample_h),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = _hex_to_rgb("#FFFFFF")
            box.line.color.rgb = _hex_to_rgb(edge)
            box.line.width = Pt(lw)
            if dash is not None:
                box.line.dash_style = dash

            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_cursor),
                Inches(legend_top + 0.02),
                Inches(sample_w * stripe_frac),
                Inches(sample_h),
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = _hex_to_rgb(stripe_color)
            stripe.line.fill.background()

            tb2 = slide.shapes.add_textbox(
                Inches(x_cursor + sample_w + 0.08),
                Inches(legend_top - 0.01),
                Inches(0.95),
                Inches(0.25),
            )
            tf2 = tb2.text_frame
            tf2.clear()
            p2 = tf2.paragraphs[0]
            p2.text = label_map.get(s, s)
            p2.font.name = font_name
            p2.font.size = Pt(10)
            p2.font.color.rgb = _hex_to_rgb("#111111")

            x_cursor += 1.18

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()
