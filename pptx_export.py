"""PowerPoint export — creates editable PPTX files.

Generates native PowerPoint shapes (not images) so users can edit
colors, text, and positions after export.
"""
import io
from datetime import date, timedelta
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from models import WorkItem, ChartConfig, STATUS_COLORS


def _hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _darken_color(hex_color: str, factor: float = 0.2) -> str:
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    r = int(r * (1 - factor))
    g = int(g * (1 - factor))
    b = int(b * (1 - factor))
    return f"#{r:02x}{g:02x}{b:02x}"


def _text_color_for_bg(hex_color: str) -> str:
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return "#FFFFFF" if luminance < 0.55 else "#1E293B"


def _lighten_color(hex_color: str, factor: float = 0.85) -> str:
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"#{r:02x}{g:02x}{b:02x}"


# ── Gantt PPTX ───────────────────────────────────────────────────────────────

def _assign_sublanes(items: List[WorkItem]):
    sorted_items = sorted(items, key=lambda it: (it.start_date, it.end_date, it.title))
    sublanes = []
    result = []
    for item in sorted_items:
        placed = False
        for lane_idx, lane_items in enumerate(sublanes):
            last_end = lane_items[-1].end_date
            if item.start_date > last_end:
                lane_items.append(item)
                result.append((item, lane_idx))
                placed = True
                break
        if not placed:
            sublanes.append([item])
            result.append((item, len(sublanes) - 1))
    return result, len(sublanes)


def export_gantt_pptx(items: List[WorkItem], config: ChartConfig) -> bytes:
    """Export Gantt chart as editable PowerPoint."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Dimensions
    slide_w = Inches(13.333)
    slide_h = Inches(7.5)
    margin = Inches(0.3)
    header_h = Inches(0.8)
    timeline_h = Inches(0.4)
    legend_h = Inches(0.4)
    sidebar_w = Inches(1.6)

    content_left = margin + sidebar_w
    content_top = margin + header_h + timeline_h
    content_w = slide_w - content_left - margin
    content_h = slide_h - content_top - legend_h - margin

    # Get categories and layout
    categories = []
    seen = set()
    for it in items:
        if it.category not in seen:
            categories.append(it.category)
            seen.add(it.category)

    from collections import defaultdict
    cats = defaultdict(list)
    for item in items:
        cats[item.category].append(item)

    # Compute sublanes per category
    cat_layout = {}
    total_sublanes = 0
    for cat in categories:
        assignments, num_sublanes = _assign_sublanes(cats[cat])
        cat_layout[cat] = {"assignments": assignments, "num_sublanes": num_sublanes}
        total_sublanes += num_sublanes

    # Date range
    chart_start = config.start_date or min(it.start_date for it in items) - timedelta(days=7)
    chart_end = config.end_date or max(it.end_date for it in items) + timedelta(days=7)
    total_days = (chart_end - chart_start).days

    # ── Header ───────────────────────────────────────────────────────────
    title_box = slide.shapes.add_textbox(margin, Inches(0.15), Inches(8), Inches(0.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = config.title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb("#0F172A")

    if config.subtitle:
        sub_box = slide.shapes.add_textbox(margin, Inches(0.55), Inches(8), Inches(0.3))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = config.subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = _hex_to_rgb("#64748B")

    # Date range text
    date_box = slide.shapes.add_textbox(Inches(9), Inches(0.3), Inches(4), Inches(0.3))
    tf3 = date_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"{chart_start.strftime('%b %d, %Y')}  →  {chart_end.strftime('%b %d, %Y')}"
    p3.font.size = Pt(10)
    p3.font.color.rgb = _hex_to_rgb("#94A3B8")
    p3.alignment = PP_ALIGN.RIGHT

    # ── Timeline header ──────────────────────────────────────────────────
    from block_renderer import _compute_time_columns
    time_columns, _ = _compute_time_columns(chart_start, chart_end)

    timeline_top = margin + header_h
    for i, (col_start, col_end, label) in enumerate(time_columns):
        x_frac_start = max(0, (col_start - chart_start).days / total_days)
        x_frac_end = min(1, (col_end - chart_start).days / total_days)

        left = content_left + int(content_w * x_frac_start)
        width = int(content_w * (x_frac_end - x_frac_start))

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, int(timeline_top), max(width, Inches(0.1)), int(timeline_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb("#334155" if i % 2 == 0 else "#1E293B")
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.color.rgb = _hex_to_rgb("#F1F5F9")
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

    # ── Swim lanes and tasks ─────────────────────────────────────────────
    if total_sublanes == 0:
        total_sublanes = 1
    sublane_h = content_h / total_sublanes
    current_y = content_top
    row_padding = Inches(0.03)

    for cat_idx, cat in enumerate(categories):
        info = cat_layout[cat]
        num_sublanes = info["num_sublanes"]
        lane_h = int(sublane_h * num_sublanes)
        color = config.get_category_color(cat, categories)
        bg_color = _lighten_color(color, 0.92)

        # Lane background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(content_left), int(current_y), int(content_w), lane_h,
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = _hex_to_rgb(bg_color)
        bg_shape.line.color.rgb = _hex_to_rgb("#E2E8F0")
        bg_shape.line.width = Pt(0.5)

        # Sidebar label
        label_shape = slide.shapes.add_textbox(
            margin, int(current_y), int(sidebar_w), lane_h,
        )
        tf = label_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cat
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = _hex_to_rgb(color)
        p.alignment = PP_ALIGN.RIGHT
        tf.paragraphs[0].space_before = Pt(0)

        # Color bar indicator
        bar_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(content_left - Inches(0.06)), int(current_y), Inches(0.06), lane_h,
        )
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = _hex_to_rgb(color)
        bar_shape.line.fill.background()

        # Tasks
        for item, sublane in info["assignments"]:
            bar_color = item.color_override if item.color_override else color
            txt_color = _text_color_for_bg(bar_color)

            x_frac_start = max(0, (item.start_date - chart_start).days / total_days)
            x_frac_end = min(1, (item.end_date - chart_start).days / total_days)
            if x_frac_end <= x_frac_start:
                x_frac_end = x_frac_start + 1 / total_days

            left = content_left + int(content_w * x_frac_start)
            width = max(int(content_w * (x_frac_end - x_frac_start)), Inches(0.15))
            top = int(current_y + sublane * sublane_h + row_padding)
            height = max(int(sublane_h - row_padding * 2), Inches(0.2))

            if item.is_milestone:
                mid_x = left + width // 2
                mid_y = top + height // 2
                size = min(Inches(0.2), height // 2)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.DIAMOND,
                    mid_x - size, mid_y - size, size * 2, size * 2,
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = _hex_to_rgb(bar_color)
                shape.line.color.rgb = _hex_to_rgb(_darken_color(bar_color))
                shape.line.width = Pt(1)
            else:
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left, top, width, height,
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = _hex_to_rgb(bar_color)
                shape.line.color.rgb = _hex_to_rgb(_darken_color(bar_color, 0.15))
                shape.line.width = Pt(0.75)

                # Task text
                tf = shape.text_frame
                tf.word_wrap = True
                tf.auto_size = None
                p = tf.paragraphs[0]
                p.text = item.title
                p.font.size = Pt(min(9, max(6, int(height / Inches(0.05)))))
                p.font.color.rgb = _hex_to_rgb(txt_color)
                p.font.bold = True
                p.alignment = PP_ALIGN.LEFT
                tf.paragraphs[0].space_before = Pt(1)
                tf.paragraphs[0].space_after = Pt(0)

        current_y += lane_h

    # Save
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ── Block Diagram PPTX ───────────────────────────────────────────────────────

def export_block_pptx(items: List[WorkItem], config: ChartConfig) -> bytes:
    """Export block diagram as editable PowerPoint."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin = Inches(0.3)
    header_h = Inches(0.8)
    timeline_h = Inches(0.35)
    legend_h = Inches(0.4)

    content_left = margin
    content_top = margin + header_h + timeline_h
    content_w = slide_w - margin * 2
    content_h = slide_h - content_top - legend_h

    categories = []
    seen = set()
    for it in items:
        if it.category not in seen:
            categories.append(it.category)
            seen.add(it.category)

    chart_start = config.start_date or min(it.start_date for it in items) - timedelta(days=3)
    chart_end = config.end_date or max(it.end_date for it in items) + timedelta(days=3)
    total_days = (chart_end - chart_start).days

    # Header
    title_box = slide.shapes.add_textbox(margin, Inches(0.15), Inches(8), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = config.title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb("#0F172A")

    if config.subtitle:
        sub_box = slide.shapes.add_textbox(margin, Inches(0.55), Inches(8), Inches(0.3))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = config.subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = _hex_to_rgb("#64748B")

    # Timeline header
    from block_renderer import _compute_time_columns, _pack_rows
    time_columns, _ = _compute_time_columns(chart_start, chart_end)
    timeline_top = margin + header_h

    for i, (col_start, col_end, label) in enumerate(time_columns):
        x_frac_start = max(0, (col_start - chart_start).days / total_days)
        x_frac_end = min(1, (col_end - chart_start).days / total_days)
        left = int(content_left + content_w * x_frac_start)
        width = max(int(content_w * (x_frac_end - x_frac_start)), Inches(0.1))

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, int(timeline_top), width, int(timeline_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb("#334155" if i % 2 == 0 else "#1E293B")
        shape.line.fill.background()

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.color.rgb = _hex_to_rgb("#F1F5F9")
        p.alignment = PP_ALIGN.CENTER

    # Pack rows
    packed_rows = _pack_rows(items, chart_start, chart_end)
    num_rows = max(1, len(packed_rows))
    row_gap = Inches(0.05)
    row_h = (content_h - row_gap * (num_rows + 1)) / num_rows

    for row_idx, row_items in enumerate(packed_rows):
        y_top = int(content_top + row_gap + row_idx * (row_h + row_gap))

        for item in row_items:
            x_frac_start = max(0, (item.start_date - chart_start).days / total_days)
            x_frac_end = min(1, (item.end_date - chart_start).days / total_days)
            if x_frac_end <= x_frac_start:
                x_frac_end = x_frac_start + 1 / total_days

            left = int(content_left + content_w * x_frac_start + Inches(0.03))
            width = max(int(content_w * (x_frac_end - x_frac_start) - Inches(0.06)), Inches(0.15))
            height = max(int(row_h), Inches(0.3))

            bar_color = config.get_category_color(item.category, categories)
            if item.color_override:
                bar_color = item.color_override
            txt_color = _text_color_for_bg(bar_color)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, y_top, width, height,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb(bar_color)
            shape.line.color.rgb = _hex_to_rgb(_darken_color(bar_color, 0.15))
            shape.line.width = Pt(0.75)

            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = None

            # Label
            if item.label:
                p = tf.paragraphs[0]
                p.text = item.label
                p.font.size = Pt(7)
                p.font.bold = True
                p.font.color.rgb = _hex_to_rgb(txt_color)
                p.font.italic = True
                p.space_before = Pt(2)
                p.space_after = Pt(0)

                p2 = tf.add_paragraph()
                p2.text = item.title
                p2.font.size = Pt(min(9, max(6, int(height / Inches(0.07)))))
                p2.font.bold = True
                p2.font.color.rgb = _hex_to_rgb(txt_color)
                p2.space_before = Pt(1)
                p2.space_after = Pt(0)
            else:
                p = tf.paragraphs[0]
                p.text = item.title
                p.font.size = Pt(min(9, max(6, int(height / Inches(0.07)))))
                p.font.bold = True
                p.font.color.rgb = _hex_to_rgb(txt_color)
                p.space_before = Pt(2)
                p.space_after = Pt(0)

            # Description
            if item.description:
                p3 = tf.add_paragraph()
                desc = item.description[:100]
                p3.text = f"• {desc}"
                p3.font.size = Pt(min(7, max(5, int(height / Inches(0.1)))))
                p3.font.color.rgb = _hex_to_rgb(txt_color)
                p3.space_before = Pt(1)

    # Legend
    legend_y = int(slide_h - legend_h)
    x_pos = int(margin)
    for cat in categories:
        color = config.get_category_color(cat, categories)

        swatch = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, legend_y + Inches(0.05), Inches(0.2), Inches(0.2),
        )
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = _hex_to_rgb(color)
        swatch.line.fill.background()

        label_box = slide.shapes.add_textbox(
            x_pos + Inches(0.25), legend_y, Inches(1.5), Inches(0.3),
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = cat
        p.font.size = Pt(9)
        p.font.color.rgb = _hex_to_rgb("#334155")
        p.font.bold = True

        x_pos += Inches(0.25 + len(cat) * 0.08 + 0.3)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ── Sequencing Diagram PPTX ──────────────────────────────────────────────────

def export_sequencing_pptx(items: List[WorkItem], config: ChartConfig) -> bytes:
    """Export sequencing diagram as editable PowerPoint."""
    from sequencing_renderer import (
        _parse_bullets, _layout_sequencing, _compute_two_row_timeline,
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin = Inches(0.3)
    header_h = Inches(0.8)
    quarter_h = Inches(0.3)
    month_h = Inches(0.25)
    sidebar_w = Inches(1.5)

    content_left = margin + sidebar_w
    content_top_y = margin + header_h + quarter_h + month_h
    content_w = slide_w - content_left - margin
    content_h = slide_h - content_top_y - margin

    categories = []
    seen = set()
    for it in items:
        if it.category not in seen:
            categories.append(it.category)
            seen.add(it.category)

    chart_start = config.start_date or min(it.start_date for it in items) - timedelta(days=3)
    chart_end = config.end_date or max(it.end_date for it in items) + timedelta(days=3)
    total_days = (chart_end - chart_start).days

    # ── Header ───────────────────────────────────────────────────────────
    title_box = slide.shapes.add_textbox(margin, Inches(0.15), Inches(8), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = config.title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb("#0F172A")

    if config.subtitle:
        sub_box = slide.shapes.add_textbox(margin, Inches(0.55), Inches(8), Inches(0.3))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = config.subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = _hex_to_rgb("#64748B")

    # Date range text
    date_box = slide.shapes.add_textbox(Inches(9), Inches(0.3), Inches(4), Inches(0.3))
    tf3 = date_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"{chart_start.strftime('%b %d, %Y')}  \u2192  {chart_end.strftime('%b %d, %Y')}"
    p3.font.size = Pt(10)
    p3.font.color.rgb = _hex_to_rgb("#94A3B8")
    p3.alignment = PP_ALIGN.RIGHT

    # ── Two-row timeline ─────────────────────────────────────────────────
    quarter_columns, month_columns = _compute_two_row_timeline(chart_start, chart_end)

    quarter_top_y = margin + header_h
    month_top_y = quarter_top_y + int(quarter_h)

    for i, (col_start, col_end, label) in enumerate(quarter_columns):
        x_frac_start = max(0, (col_start - chart_start).days / total_days)
        x_frac_end = min(1, (col_end - chart_start).days / total_days)
        left = int(content_left + content_w * x_frac_start)
        width = max(int(content_w * (x_frac_end - x_frac_start)), Inches(0.1))

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, int(quarter_top_y), width, int(quarter_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb("#334155" if i % 2 == 0 else "#1E293B")
        shape.line.fill.background()

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = _hex_to_rgb("#F1F5F9")
        p.alignment = PP_ALIGN.CENTER

    for i, (col_start, col_end, label) in enumerate(month_columns):
        x_frac_start = max(0, (col_start - chart_start).days / total_days)
        x_frac_end = min(1, (col_end - chart_start).days / total_days)
        left = int(content_left + content_w * x_frac_start)
        width = max(int(content_w * (x_frac_end - x_frac_start)), Inches(0.1))

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, int(month_top_y), width, int(month_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb("#546378" if i % 2 == 0 else "#475569")
        shape.line.fill.background()

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.color.rgb = _hex_to_rgb("#E2E8F0")
        p.alignment = PP_ALIGN.CENTER

    # ── Left-side legend ─────────────────────────────────────────────────
    if config.show_legend:
        legend_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(margin), int(content_top_y), int(sidebar_w - Inches(0.1)), int(content_h),
        )
        legend_bg.fill.solid()
        legend_bg.fill.fore_color.rgb = _hex_to_rgb("#F1F5F9")
        legend_bg.line.color.rgb = _hex_to_rgb("#E2E8F0")
        legend_bg.line.width = Pt(0.5)

        # Legend title
        legend_title = slide.shapes.add_textbox(
            int(margin + Inches(0.1)), int(content_top_y + Inches(0.1)),
            int(sidebar_w - Inches(0.3)), Inches(0.25),
        )
        tf = legend_title.text_frame
        p = tf.paragraphs[0]
        p.text = "LEGEND"
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = _hex_to_rgb("#64748B")

        # Category entries
        cat_y = content_top_y + Inches(0.4)
        for cat in categories:
            color = config.get_category_color(cat, categories)

            swatch = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(margin + Inches(0.1)), int(cat_y),
                Inches(0.2), Inches(0.2),
            )
            swatch.fill.solid()
            swatch.fill.fore_color.rgb = _hex_to_rgb(color)
            swatch.line.fill.background()

            label_box = slide.shapes.add_textbox(
                int(margin + Inches(0.35)), int(cat_y - Inches(0.02)),
                int(sidebar_w - Inches(0.5)), Inches(0.25),
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = cat
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb("#334155")

            cat_y += Inches(0.35)

    # ── Blocks ───────────────────────────────────────────────────────────
    content_w_inches = content_w / Inches(1)
    content_h_inches = content_h / Inches(1)
    base_line_h_inches = 8.5 * 0.016
    base_line_h = base_line_h_inches / content_h_inches

    layouts = _layout_sequencing(
        items, chart_start, chart_end,
        content_w_inches, content_h_inches,
        base_line_h, 8.5,
    )

    # Scale layouts to fit if they overflow
    if layouts:
        min_y = min(l["y_bot"] for l in layouts)
        if min_y < 0:
            scale = 1.0 / (1.0 + abs(min_y) + 0.02)
            for l in layouts:
                l["y_top"] = 1.0 - (1.0 - l["y_top"]) / (1.0 + abs(min_y) + 0.02)
                l["y_bot"] = l["y_top"] - l["height"] * scale
                l["height"] = l["height"] * scale

    for layout in layouts:
        item = layout["item"]
        x_s = layout["x_start"]
        x_e = layout["x_end"]
        y_top_frac = layout["y_top"]
        y_bot_frac = layout["y_bot"]

        bar_color = config.get_category_color(item.category, categories)
        if item.color_override:
            bar_color = item.color_override
        fill_color = _lighten_color(bar_color, 0.55)
        txt_color = _text_color_for_bg(fill_color)
        border = _darken_color(bar_color, 0.1)

        # Draw all segment rectangles for this item
        rects = layout.get("rects", [])
        for rect in rects:
            rx_s = rect["x_start"]
            rx_e = rect["x_end"]
            ry_top = rect["y_top"]
            ry_bot = rect["y_bot"]
            r_left = int(content_left + content_w * rx_s)
            r_width = max(int(content_w * (rx_e - rx_s)), Inches(0.1))
            r_top = int(content_top_y + content_h * (1.0 - ry_top))
            r_height = max(int(content_h * (ry_top - ry_bot)), Inches(0.1))

            # Skip if this is the primary rect (drawn below with text)
            if (abs(rx_s - x_s) < 0.002 and abs(rx_e - x_e) < 0.002 and
                abs(ry_top - y_top_frac) < 0.002):
                continue

            extra_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                r_left, r_top, r_width, r_height,
            )
            extra_shape.fill.solid()
            extra_shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
            extra_shape.line.color.rgb = _hex_to_rgb(border)
            extra_shape.line.width = Pt(0.5)

        # Primary rectangle with text
        left = int(content_left + content_w * x_s)
        width = max(int(content_w * (x_e - x_s)), Inches(0.2))
        top = int(content_top_y + content_h * (1.0 - y_top_frac))
        height = max(int(content_h * (y_top_frac - y_bot_frac)), Inches(0.3))

        if item.is_milestone:
            mid_x = left + width // 2
            mid_y = top + height // 2
            size = min(Inches(0.2), height // 2)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                mid_x - size, mid_y - size, size * 2, size * 2,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb(bar_color)
            shape.line.color.rgb = _hex_to_rgb(border)
            shape.line.width = Pt(1)
            continue

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
        shape.line.color.rgb = _hex_to_rgb(border)
        shape.line.width = Pt(0.75)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        font_sz = min(9, max(6, int(height / Inches(0.06))))

        # Label
        if item.label:
            p = tf.paragraphs[0]
            p.text = item.label
            p.font.size = Pt(max(6, font_sz - 1))
            p.font.bold = True
            p.font.italic = True
            p.font.color.rgb = _hex_to_rgb("#475569")
            p.space_before = Pt(2)
            p.space_after = Pt(0)

            p2 = tf.add_paragraph()
            title_text = item.title
            if item.status == "done":
                title_text += "  \u2713"
            p2.text = title_text
            p2.font.size = Pt(font_sz)
            p2.font.bold = True
            p2.font.color.rgb = _hex_to_rgb("#1E293B")
            p2.space_before = Pt(1)
            p2.space_after = Pt(0)
        else:
            p = tf.paragraphs[0]
            title_text = item.title
            if item.status == "done":
                title_text += "  \u2713"
            p.text = title_text
            p.font.size = Pt(font_sz)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb("#1E293B")
            p.space_before = Pt(2)
            p.space_after = Pt(0)

        # Bullet points
        bullets = _parse_bullets(item.description)
        bullet_fs = max(5, font_sz - 2)
        for bullet in bullets:
            bp = tf.add_paragraph()
            bp.text = f"\u2022 {bullet[:80]}"
            bp.font.size = Pt(bullet_fs)
            bp.font.color.rgb = _hex_to_rgb("#334155")
            bp.space_before = Pt(0)
            bp.space_after = Pt(0)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
